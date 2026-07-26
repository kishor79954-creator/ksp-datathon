from pptx import Presentation

prs = Presentation(r"C:\Users\Kishore\Downloads\KSP Datathon 2026 _ Prototype Submission Template.pptx")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        print(f"Shape {j}: name='{shape.name}', type={shape.shape_type}")
        if shape.is_placeholder:
            print(f"  Placeholder: type={shape.placeholder_format.type}")
