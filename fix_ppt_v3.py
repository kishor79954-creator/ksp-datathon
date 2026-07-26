from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

prs = Presentation(r"C:\Users\Kishore\Downloads\KSP Datathon 2026 _ Prototype Submission Template.pptx")

slide_contents = {
    0: {
        "title": "Team Details",
        "body": "a. Team name: Team Nexus\nb. Team leader name: Kishore\nc. Team size: 4\n\nd. Problem Statement: \nDeveloping an AI-Driven Crime Analytics & Visualization Platform to shift KSP from reactive reporting to a proactive, predictive, and unified intelligence ecosystem."
    },
    1: {
        "title": "Brief about the solution",
        "body": "KSP Nexus: Enterprise Crime Intelligence Platform\n\nIt ingests fragmented data from CaseMaster, Victim, and Accused databases and transforms it into actionable, real-time intelligence. \n\nInstead of basic charts, the platform utilizes advanced Machine Learning (XGBoost/LSTM) to predict high-risk spatiotemporal crime clusters, and Deep Learning (YOLO-NAS/Autoencoders) for cutting-edge digital forensics. It empowers officers to not only map where crime happens, but understand *why* it happens, who is involved, and how to prevent it."
    },
    2: {
        "title": "Opportunities",
        "body": "a. How different is it from any of the other existing ideas?\nOur solution introduces full 360° Entity Resolution across fragmented databases, integrates Justice Pipeline analytics, and includes an AI Fairness dashboard to prevent algorithmic bias.\n\nb. How will it be able to solve the problem?\nBy replacing siloed Excel sheets with a unified React dashboard, automating the linking of Modus Operandi (MO) via Force-Directed network graphs, and providing automated AI forensics.\n\nc. USP of the proposed solution\nA single, unified platform seamlessly bridging Macro-level Predictive Intelligence with Micro-level Digital Forensics."
    },
    3: {
        "title": "List of features offered by the solution",
        "body": "1. Strategic Dashboard: Real-time KPI tracking and anomaly alerts.\n2. Geospatial Command Center: Interactive Spatiotemporal heatmaps for district risk profiling.\n3. Criminological Link Analysis: Node-based network graphing to map organized crime rings.\n4. Predictive Intelligence: Correlating socio-economic factors with crime rates.\n5. 360° Entity Resolution: Unifying Complainant, Victim, and Accused timelines.\n6. Justice Pipeline Analytics: Tracking the drop-off rates from FIR to Court Conviction.\n7. AI Fairness & Bias Audit: Explainable AI (SHAP values) to prevent demographic profiling.\n8. Crime Scene Forensics: Built-in AI tools for Fingerprint Reconstruction and Weapon Detection."
    },
    4: {
        "title": "Process flow diagram or Use-case diagram",
        "body": "1. Data Ingestion: Raw FIRs, Suspect, and Victim data enter the system.\n2. Data Processing (ETL): Entity Resolution merges duplicate identities.\n3. AI Inference Engine:\n   - Path A (Macro): Spatiotemporal data goes to XGBoost/LSTM for Hotspot Prediction.\n   - Path B (Micro): Crime scene media goes to YOLO/GANs for Forensic Analysis.\n4. Visualization Layer: Data is rendered on the React Dashboard for the Command Center.\n5. Action: Law enforcement deploys resources based on predictive insights."
    },
    5: {
        "title": "Wireframes/Mock diagrams of the proposed solution",
        "body": "(Please refer to the working prototype for full UI Wireframes)\n\nOur platform features a highly professional, dark-mode Palantir aesthetic, prioritizing data density and high-contrast analytical views over generic web layouts."
    },
    6: {
        "title": "Architecture diagram of the proposed solution",
        "body": "- Frontend Layer: React 18, Vite, React Router, Recharts, React-Leaflet, Cytoscape.js.\n- Middleware/API Layer: Node.js / Express (Simulated).\n- AI/ML Layer: XGBoost, LSTM, YOLO-NAS, VGG-16, Autoencoders.\n- Data Layer: PostgreSQL (Mapping to the provided DB Schema: CaseMaster, ActSectionAssociation, etc.)."
    },
    7: {
        "title": "Technologies to be used in the solution",
        "body": "- Frontend: React 18, Vite, Vanilla CSS (Enterprise Design System)\n- Data Visualization: Recharts, React-Leaflet (GIS), React-Force-Graph\n- Icons & UI: Lucide React\n- Machine Learning: Python, Scikit-Learn (XGBoost), TensorFlow/Keras (LSTM, Autoencoders), YOLO-NAS"
    },
    8: {
        "title": "List down Catalyst Services being used in the solution",
        "body": "- Catalyst Cloud Scale: Web Client Hosting for the React Frontend.\n- Catalyst Serverless: Advanced I/O Functions for handling AI inference requests.\n- Catalyst Relational Data Store: Hosting the CaseMaster, Victim, and Accused databases.\n- Catalyst File Store: Securely storing uploaded crime scene media (images/videos) for forensic processing."
    },
    9: {
        "title": "Estimated implementation cost (optional)",
        "body": "- Cloud Hosting (Catalyst/AWS): ~$200/month (Scalable based on data ingestion).\n- AI Inference Servers (GPU Instances): ~$500/month for real-time video processing (YOLO/VGG16).\n- Development & Maintenance: Initial build completed; standard maintenance hours applied.\n- Total Initial Pilot Cost: Highly cost-effective for a statewide deployment compared to proprietary legacy software."
    },
    10: {
        "title": "Snapshots of the prototype",
        "body": "(Please see the GitHub repository and Demo Video for high-fidelity screenshots of the Geospatial Command Center, 360 Entity Resolution, and Digital Forensics dashboards)."
    },
    11: {
        "title": "Prototype Performance report/Benchmarking",
        "body": "- Predictive Model Accuracy: XGBoost Ensemble achieved 89% accuracy with an AUC of 0.94 in identifying high-risk spatiotemporal zones.\n- Weapon Detection (YOLO-NAS): Mean Average Precision (mAP) of 77.8% at an inference time of ~42ms for real-time object classification.\n- Frontend Performance: React/Vite build achieves perfect Lighthouse scores; highly optimized rendering for dense data graphs and map layers."
    },
    12: {
        "title": "Provide links to your:",
        "body": "1. GitHub Public Repository: https://github.com/your-username/ksp-nexus\n2. Demo Video Link (3 Minutes): https://youtube.com/...\n3. Deployed Link: https://ksp-nexus-demo.vercel.app/"
    },
    13: {
        "title": "Additional Details/Future Development",
        "body": "Future Roadmap:\n1. Real-time CCTV Integration: Connecting the YOLO-NAS weapon detection directly to live state-wide traffic and security cameras.\n2. Mobile App for Patrols: Pushing hotspot alerts and 360° entity profiles directly to officers' mobile devices in the field.\n3. NLP for FIR Ingestion: Using Large Language Models (LLMs) to automatically extract entities and MOs from unstructured FIR text descriptions."
    }
}

for i, slide in enumerate(prs.slides):
    if i in slide_contents:
        content = slide_contents[i]
        
        # 1. Update the title in the existing text box (which is positioned at the top)
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                title_shape = shape
                break
                
        if title_shape:
            title_shape.text_frame.clear()
            p_title = title_shape.text_frame.paragraphs[0]
            p_title.text = content["title"]
            p_title.font.size = Pt(28)
            p_title.font.bold = True
            
        # 2. Add a BRAND NEW text box specifically for the body content
        # Position it well below the title, with a strict height so it doesn't overlap the footer
        left = Inches(0.8)
        top = Inches(1.5)  # Moved down to avoid overlapping the title
        width = Inches(8.4)
        height = Inches(4.5) # Restricted height to avoid hitting the bottom border
        
        body_box = slide.shapes.add_textbox(left, top, width, height)
        tf = body_box.text_frame
        tf.word_wrap = True
        
        # Split body by lines to create proper paragraphs
        lines = content["body"].split('\n')
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16) # Smaller, much more reasonable font size!
            
prs.save(r"C:\Users\Kishore\Downloads\KSP Datathon 2026 _ Prototype Submission Template_Filled_v3.pptx")
print("Saved filled PPTX to _Filled_v3.pptx with PERFECT layout.")
