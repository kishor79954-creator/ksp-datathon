from pptx import Presentation

prs = Presentation(r"C:\Users\Kishore\Downloads\KSP Datathon 2026 _ Prototype Submission Template.pptx")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text.replace('\n', ' ')
        print(f"Shape {j}: {text[:100]}")
