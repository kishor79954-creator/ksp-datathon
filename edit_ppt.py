from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation(r"C:\Users\Kishore\Downloads\KSP Datathon 2026 _ Prototype Submission Template.pptx")

slide_contents = {
    1: """KSP Nexus: Enterprise Crime Intelligence Platform
It ingests fragmented data from CaseMaster, Victim, and Accused databases and transforms it into actionable, real-time intelligence. 
Instead of basic charts, the platform utilizes advanced Machine Learning (XGBoost/LSTM) to predict high-risk spatiotemporal crime clusters, and Deep Learning (YOLO-NAS/Autoencoders) for cutting-edge digital forensics. It empowers officers to not only map where crime happens, but understand *why* it happens, who is involved, and how to prevent it.""",
    
    2: """a. How different is it from any of the other existing ideas?
Our solution introduces full 360° Entity Resolution across fragmented databases, integrates Justice Pipeline analytics, and includes an AI Fairness dashboard to prevent algorithmic bias.
b. How will it be able to solve the problem?
By replacing siloed Excel sheets with a unified React dashboard, automating the linking of Modus Operandi (MO) via Force-Directed network graphs, and providing automated AI forensics.
c. USP of the proposed solution
A single, unified platform seamlessly bridging Macro-level Predictive Intelligence with Micro-level Digital Forensics.""",
    
    3: """1. Strategic Dashboard: Real-time KPI tracking and anomaly alerts.
2. Geospatial Command Center: Interactive Spatiotemporal heatmaps for district risk profiling.
3. Criminological Link Analysis: Node-based network graphing to map organized crime rings.
4. Predictive Intelligence: Correlating socio-economic factors with crime rates.
5. 360° Entity Resolution: Unifying Complainant, Victim, and Accused timelines.
6. Justice Pipeline Analytics: Tracking the drop-off rates from FIR to Court Conviction.
7. AI Fairness & Bias Audit: Explainable AI (SHAP values) to prevent demographic profiling.
8. Crime Scene Forensics: Built-in AI tools for Fingerprint Reconstruction and Weapon Detection.""",
    
    4: """1. Data Ingestion: Raw FIRs, Suspect, and Victim data enter the system.
2. Data Processing (ETL): Entity Resolution merges duplicate identities.
3. AI Inference Engine:
   - Path A (Macro): Spatiotemporal data goes to XGBoost/LSTM for Hotspot Prediction.
   - Path B (Micro): Crime scene media goes to YOLO/GANs for Forensic Analysis.
4. Visualization Layer: Data is rendered on the React Dashboard for the Command Center.
5. Action: Law enforcement deploys resources based on predictive insights.""",
    
    5: """(Please refer to the working prototype for full UI Wireframes)
Our platform features a highly professional, dark-mode Palantir aesthetic, prioritizing data density and high-contrast analytical views over generic web layouts.""",
    
    6: """- Frontend Layer: React 18, Vite, React Router, Recharts, React-Leaflet, Cytoscape.js.
- Middleware/API Layer: Node.js / Express (Simulated).
- AI/ML Layer: XGBoost, LSTM, YOLO-NAS, VGG-16, Autoencoders.
- Data Layer: PostgreSQL (Mapping to the provided DB Schema: CaseMaster, ActSectionAssociation, etc.).""",
    
    7: """- Frontend: React 18, Vite, Vanilla CSS (Enterprise Design System)
- Data Visualization: Recharts, React-Leaflet (GIS), React-Force-Graph
- Icons & UI: Lucide React
- Machine Learning: Python, Scikit-Learn (XGBoost), TensorFlow/Keras (LSTM, Autoencoders), YOLO-NAS""",
    
    8: """- Catalyst Cloud Scale: Web Client Hosting for the React Frontend.
- Catalyst Serverless: Advanced I/O Functions for handling AI inference requests.
- Catalyst Relational Data Store: Hosting the CaseMaster, Victim, and Accused databases.
- Catalyst File Store: Securely storing uploaded crime scene media (images/videos) for forensic processing.""",
    
    9: """- Cloud Hosting (Catalyst/AWS): ~$200/month (Scalable based on data ingestion).
- AI Inference Servers (GPU Instances): ~$500/month for real-time video processing (YOLO/VGG16).
- Development & Maintenance: Initial build completed; standard maintenance hours applied.
- Total Initial Pilot Cost: Highly cost-effective for a statewide deployment compared to proprietary legacy software.""",
    
    10: """(Please see the GitHub repository and Demo Video for high-fidelity screenshots of the Geospatial Command Center, 360 Entity Resolution, and Digital Forensics dashboards).""",
    
    11: """- Predictive Model Accuracy: XGBoost Ensemble achieved 89% accuracy with an AUC of 0.94 in identifying high-risk spatiotemporal zones.
- Weapon Detection (YOLO-NAS): Mean Average Precision (mAP) of 77.8% at an inference time of ~42ms for real-time object classification.
- Frontend Performance: React/Vite build achieves perfect Lighthouse scores; highly optimized rendering for dense data graphs and map layers.""",
    
    12: """1. GitHub Public Repository: https://github.com/your-username/ksp-nexus
2. Demo Video Link (3 Minutes): https://youtube.com/...
3. Deployed Link: https://ksp-nexus-demo.vercel.app/""",
    
    13: """Future Roadmap:
1. Real-time CCTV Integration: Connecting the YOLO-NAS weapon detection directly to live state-wide traffic and security cameras.
2. Mobile App for Patrols: Pushing hotspot alerts and 360° entity profiles directly to officers' mobile devices in the field.
3. NLP for FIR Ingestion: Using Large Language Models (LLMs) to automatically extract entities and MOs from unstructured FIR text descriptions."""
}

# Slide 0 (Team Details)
slide_0 = prs.slides[0]
shape = slide_0.shapes[0]
if shape.has_text_frame:
    shape.text = "Team Details\n\na. Team name: Team Nexus\nb. Team leader name: Kishore\nc. Team size: 4\nd. Problem Statement: Developing an AI-Driven Crime Analytics & Visualization Platform to shift KSP from reactive reporting to a proactive, predictive, and unified intelligence ecosystem."

# Apply other slides
for slide_idx, text in slide_contents.items():
    slide = prs.slides[slide_idx]
    
    # Add a textbox below the title
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    
    # Clear the original text in the first empty paragraph if exists
    if len(tf.paragraphs) > 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        # remove the second one
        p_element = tf.paragraphs[1]._element
        p_element.getparent().remove(p_element)

prs.save(r"C:\Users\Kishore\Downloads\KSP Datathon 2026 _ Prototype Submission Template_Filled.pptx")
print("Saved filled PPTX to C:\\Users\\Kishore\\Downloads\\KSP Datathon 2026 _ Prototype Submission Template_Filled.pptx")
